from collections import defaultdict
from django.shortcuts import render, redirect, get_object_or_404
from django.urls import reverse
from core.models import Event, Ssp, EventSsp
# from templatetags.custom_filters 
from event_creation.templatetags.custom_filters import get_next_category,get_previous_category
#PPT Creation import
import os
from django.http import HttpResponse,JsonResponse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from django.conf import settings

categories = ['Furniture', 'Decor', 'Stage', 'Setup', 'Entrance', 'Services', 'Audio-Visual', 'Stalls', 'Lighting', 'Miscellaneous']

def dashboard_ec(request,event_id):
    print(event_id)
    event = get_object_or_404(Event, event_id=event_id)
    return render(request, 'dashboard_ec.html', {'event': event})


def start_flow(request, event_id):
    event = get_object_or_404(Event, event_id=event_id)
    # Start the flow by redirecting to the first category
    return redirect('event_creation:category_flow', event_id=event.event_id, category_name=categories[0])

def summary(request, event_id):
    event = get_object_or_404(Event, event_id=event_id)
    linked_ssps = EventSsp.objects.filter(event=event).select_related('ssp')

    context = {
        'event': event,
        'linked_ssps': linked_ssps,
    }
    return render(request, 'summary.html', context)


def quotation(request, event_id):
    event = get_object_or_404(Event, event_id=event_id)
    linked_ssps = EventSsp.objects.filter(event=event).select_related('ssp')

    # Group SSPs by category
    categories = defaultdict(list)
    category_totals = {}  # To store total per category
    grand_total = 0  # To store the total across all categories

    for linked_ssp in linked_ssps:
        category_name = linked_ssp.ssp.item_type  # Assuming 'category' is a field in SSP model
        categories[category_name].append(linked_ssp)
        
        # Calculate category total
        total_price = linked_ssp.quantity * linked_ssp.ssp.rate_per_event
        category_totals[category_name] = category_totals.get(category_name, 0) + total_price
        grand_total += total_price  # Add to grand total

    context = {
        'event': event,
        'categories': dict(categories),  # Convert defaultdict to normal dict
        'category_totals': category_totals,
        'grand_total': grand_total,
    }
    return render(request, 'quotation.html', context)


def category_flow(request, event_id, category_name):
    event = get_object_or_404(Event, event_id=event_id)
    ssps = Ssp.objects.filter(item_type=category_name)  # Fetch products for this category
    linked_ssps = {event_ssp.ssp.item_id: event_ssp for event_ssp in EventSsp.objects.filter(event=event)}

    if request.method == "POST":
        ssp_ids = request.POST.getlist("ssp_ids[]")
        quantities = request.POST.getlist("quantities[]")

        for ssp_id, quantity in zip(ssp_ids, quantities):
            quantity = int(quantity)
            ssp = Ssp.objects.get(item_id=ssp_id)

            if quantity > 0:
                event_ssp, created = EventSsp.objects.get_or_create(event=event, ssp=ssp)
                event_ssp.quantity = quantity
                event_ssp.save()
            else:
                EventSsp.objects.filter(event=event, ssp=ssp).delete()

        # Handle navigation
        navigation = request.POST.get("navigation")
        if navigation == "previous":
            return redirect("event_creation:category_flow", event_id=event.event_id, category_name=get_previous_category(category_name))
        elif navigation == "next":
            return redirect("event_creation:category_flow", event_id=event.event_id, category_name=get_next_category(category_name))
        else:  # Finish
            return redirect("event_creation:summary", event_id=event.event_id)

    context = {
        "event": event,
        "ssps": ssps,
        "category_name": category_name,
        "categories": categories,
        "linked_ssps": linked_ssps,
    }
    return render(request, "category_flow.html", context)



def generate_ppt(request, event_id):
    event = get_object_or_404(Event, event_id=event_id)
    linked_ssps = EventSsp.objects.filter(event=event).select_related('ssp')

    prs = Presentation()

    # 🎯 **1️⃣ Cover Slide with Event Details**
    cover_slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank slide for more control
    cover_slide.background.fill.solid()
    cover_slide.background.fill.fore_color.rgb = RGBColor(154, 230, 174)  #  background

    title_box = cover_slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = f"{event.event_name}  \n {event.group}"
    title_frame.paragraphs[0].font.size = Pt(40)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # 

    subtitle_box = cover_slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = f"PPT for Event ID: {event.event_id}\nDate: {event.date}"
    subtitle_frame.paragraphs[0].font.size = Pt(30)
    subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # 

    # 🎯 **2️⃣ Group SSPs by Category**
    categories = defaultdict(list)
    for linked_ssp in linked_ssps:
        categories[linked_ssp.ssp.item_type].append(linked_ssp)

    # 🎯 **3️⃣ Loop Through Categories**
    for category, items in categories.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout for custom design
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(154, 230, 174)  # Light gray background

        # **Category Title on Top**
        cat_title_box = slide.shapes.add_textbox(Inches(3), Inches(0.2), Inches(9), Inches(0.8))
        cat_title_frame = cat_title_box.text_frame
        
        cat_title_frame.text = f"Category: {category}"
        cat_title_frame.paragraphs[0].font.size = Pt(24)
        cat_title_frame.paragraphs[0].font.bold = True
        cat_title_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)  # 

        # 🎯 **4️⃣ Arrange Products in a 3x3 Grid (9 Items Per Slide)**
        max_cols = 3
        max_rows = 3
        item_width = 2.5  # Each card width
        item_height = 3  # Each card height
        padding_x = 0.5  # X-axis padding
        padding_y = 1  # Y-axis padding

        row, col = 0, 0

        for index, linked_ssp in enumerate(items):
            ssp = linked_ssp.ssp

            # If we exceed 9 items, create a new slide
            if index % (max_cols * max_rows) == 0 and index > 0:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                slide.background.fill.solid()
                slide.background.fill.fore_color.rgb = RGBColor(154, 230, 174)  # Light gray
                row, col = 0, 0  # Reset positions

            x = Inches(padding_x + col * (item_width + 0.5))
            y = Inches(padding_y + row * (item_height + 0.5))

            # **Draw Card**
            card = slide.shapes.add_shape(1, x, y, Inches(item_width), Inches(item_height))
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor(200, 230, 255)  # Light blue

            # **Add Image**
            if ssp.image:
                image_path = os.path.join(settings.MEDIA_ROOT, ssp.image.name)
                if os.path.exists(image_path):
                    slide.shapes.add_picture(image_path, x + Inches(0.2), y + Inches(0.2), width=Inches(2))

            # **Add Name Below Image**
            name_box = slide.shapes.add_textbox(x, y + Inches(2.2), Inches(item_width), Inches(0.5))
            name_frame = name_box.text_frame
            name_frame.text = ssp.name
            name_frame.paragraphs[0].font.size = Pt(16)
            name_frame.paragraphs[0].font.bold = True
            name_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

            # **Add Rent and Quantity Below Name**
            info_box = slide.shapes.add_textbox(x, y + Inches(2.6), Inches(item_width), Inches(0.5))
            info_frame = info_box.text_frame
            info_frame.text = f"₹{ssp.rate_per_event} x {linked_ssp.quantity}"
            info_frame.paragraphs[0].font.size = Pt(12)
            info_frame.paragraphs[0].font.color.rgb = RGBColor(50, 50, 50)

            # Update row & column positions
            col += 1
            if col >= max_cols:
                col = 0
                row += 1

    # 🎯 **5️⃣ Summary Slide**
    total_cost = sum(linked_ssp.quantity * linked_ssp.ssp.rate_per_event for linked_ssp in linked_ssps)
    summary_slide = prs.slides.add_slide(prs.slide_layouts[5])
    summary_slide.background.fill.solid()
    summary_slide.background.fill.fore_color.rgb = RGBColor(154, 230, 174)  # background

    summary_title = summary_slide.shapes.add_textbox(Inches(3), Inches(2), Inches(8), Inches(1))
    summary_title.text_frame.text = "Quotation Summary"
    summary_title.text_frame.paragraphs[0].font.size = Pt(32)
    summary_title.text_frame.paragraphs[0].font.bold = True
    summary_title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)

    summary_text = summary_slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(8), Inches(1))
    summary_text.text_frame.text = f"Total Cost: ₹{total_cost}"
    summary_text.text_frame.paragraphs[0].font.size = Pt(24)
    summary_text.text_frame.paragraphs[0].font.color.rgb = RGBColor(0,0,0)

    # 🎯 **6️⃣ Save and Serve PPT**
    ppt_filename = f"SSP_PPT_{event.event_id}.pptx"
    ppt_path = os.path.join(settings.MEDIA_ROOT, ppt_filename)
    prs.save(ppt_path)

    with open(ppt_path, "rb") as ppt_file:
        response = HttpResponse(ppt_file.read(), content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation")
        response['Content-Disposition'] = f'attachment; filename="{ppt_filename}"'
        return response
